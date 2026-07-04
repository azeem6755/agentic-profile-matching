import os
import datetime
from pypdf import PdfReader
import docx2txt
from pptx import Presentation

# Default sandbox is the project root (where this file is located)
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

def check_sandbox(filepath: str, sandbox_dir: str = PROJECT_ROOT) -> str:
    """
    Resolves the absolute path of the target filepath and verifies it lies
    within the absolute path of the sandbox_dir.
    Raises PermissionError if the target path escapes the sandbox.
    """
    abs_sandbox = os.path.abspath(sandbox_dir)
    # Target path might be relative to sandbox or project root
    if not os.path.isabs(filepath):
        abs_target = os.path.abspath(os.path.join(abs_sandbox, filepath))
    else:
        abs_target = os.path.abspath(filepath)
        
    # Check traversal
    # We add trailing path separator to avoid partial prefix matching (e.g., /app vs /app-other)
    sandbox_prefix = abs_sandbox if abs_sandbox.endswith(os.sep) else abs_sandbox + os.sep
    
    # We also allow the exact sandbox directory itself
    if abs_target != abs_sandbox and not abs_target.startswith(sandbox_prefix):
        raise PermissionError(f"Access Denied: Path '{filepath}' resolves to '{abs_target}', which is outside the sandbox '{abs_sandbox}'.")
    
    return abs_target

def read_file(filepath: str) -> dict:
    """
    Reads resume files (PDF, TXT, DOCX), extracts text content,
    and returns a structured response with content and metadata.
    Handles errors gracefully.
    """
    try:
        abs_path = check_sandbox(filepath)
        if not os.path.exists(abs_path):
            return {
                "status": "error",
                "message": f"File '{filepath}' not found."
            }
        
        if os.path.isdir(abs_path):
            return {
                "status": "error",
                "message": f"Path '{filepath}' is a directory, not a file."
            }

        file_size = os.path.getsize(abs_path)
        mod_time = datetime.datetime.fromtimestamp(os.path.getmtime(abs_path)).isoformat()
        
        ext = os.path.splitext(abs_path)[1].lower()
        content = ""
        
        if ext == ".txt":
            # Plain Text
            with open(abs_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        elif ext == ".pdf":
            # PDF Extraction using pypdf
            reader = PdfReader(abs_path)
            pages_text = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            content = "\n".join(pages_text)
        elif ext == ".docx":
            # DOCX Extraction using docx2txt
            content = docx2txt.process(abs_path)
        elif ext == ".pptx":
            # PPTX Extraction using python-pptx
            prs = Presentation(abs_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            content = "\n".join(text_runs)
        else:
            return {
                "status": "error",
                "message": f"Unsupported file extension '{ext}'. Only .pdf, .docx, .pptx, and .txt are supported."
            }
            
        word_count = len(content.split())
        char_count = len(content)

        return {
            "status": "success",
            "content": content,
            "metadata": {
                "filepath": filepath,
                "format": ext[1:],
                "size_bytes": file_size,
                "modified_time": mod_time,
                "word_count": word_count,
                "char_count": char_count
            }
        }
        
    except PermissionError as pe:
        return {
            "status": "error",
            "message": str(pe)
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"An error occurred while reading file: {str(e)}"
        }

def list_files(directory: str, extension: str = None) -> list:
    """
    Lists all files in a directory, optionally filtering by extension.
    Returns metadata for each file.
    """
    try:
        abs_dir = check_sandbox(directory)
        if not os.path.exists(abs_dir):
            return [{"status": "error", "message": f"Directory '{directory}' not found."}]
            
        if not os.path.isdir(abs_dir):
            return [{"status": "error", "message": f"Path '{directory}' is a file, not a directory."}]
            
        files_metadata = []
        # Standardize extension prefix
        if extension:
            extension = extension.lower()
            if not extension.startswith("."):
                extension = "." + extension
                
        for root, dirs, files in os.walk(abs_dir):
            # Resolve root to absolute path
            abs_root = os.path.abspath(root)
            # Ensure subdirectory is also within the sandbox
            try:
                check_sandbox(abs_root)
            except PermissionError:
                continue # Skip outside folders
                
            for file in files:
                abs_file_path = os.path.join(abs_root, file)
                if extension and not file.lower().endswith(extension):
                    continue
                    
                # Relative path from project root or target directory for display
                rel_path = os.path.relpath(abs_file_path, PROJECT_ROOT)
                
                try:
                    size = os.path.getsize(abs_file_path)
                    mod_time = datetime.datetime.fromtimestamp(os.path.getmtime(abs_file_path)).isoformat()
                    files_metadata.append({
                        "name": file,
                        "relative_path": rel_path,
                        "size_bytes": size,
                        "modified_time": mod_time
                    })
                except OSError:
                    # Skip files that can't be read or accessed
                    continue
                    
        return files_metadata
        
    except PermissionError as pe:
        return [{"status": "error", "message": str(pe)}]
    except Exception as e:
        return [{"status": "error", "message": f"An error occurred while listing files: {str(e)}"}]

def write_file(filepath: str, content: str) -> dict:
    """
    Writes content to a file, creating directories if needed.
    Returns success/failure status.
    """
    try:
        abs_path = check_sandbox(filepath)
        
        # Prevent rewriting system files or key python source files
        basename = os.path.basename(abs_path)
        if basename in ["fs_tools.py", "llm_file_assistant.py", "requirements.txt"]:
            return {
                "status": "error",
                "message": f"Write Access Denied: Modifying source file '{basename}' is not permitted."
            }
            
        # Ensure directories exist
        parent_dir = os.path.dirname(abs_path)
        os.makedirs(parent_dir, exist_ok=True)
        
        with open(abs_path, "w", encoding="utf-8") as f:
            f.write(content)
            
        size = os.path.getsize(abs_path)
        return {
            "status": "success",
            "filepath": filepath,
            "size_bytes_written": size,
            "message": f"Successfully wrote {size} bytes to '{filepath}'."
        }
        
    except PermissionError as pe:
        return {
            "status": "error",
            "message": str(pe)
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"An error occurred while writing file: {str(e)}"
        }

def search_in_file(filepath: str, keyword: str) -> dict:
    """
    Searches for a case-insensitive keyword in file content.
    Returns matches with context (surrounding text).
    """
    try:
        abs_path = check_sandbox(filepath)
        
        # Read content using read_file to handle text/pdf/docx format parsing
        read_res = read_file(filepath)
        if read_res["status"] == "error":
            return read_res
            
        content = read_res["content"]
        lines = content.splitlines()
        matches = []
        keyword_lower = keyword.lower()
        
        for idx, line in enumerate(lines):
            if keyword_lower in line.lower():
                matches.append({
                    "line_number": idx + 1,
                    "line_content": line.strip()
                })
                
        return {
            "status": "success",
            "filepath": filepath,
            "keyword": keyword,
            "occurrences": len(matches),
            "matches": matches
        }
        
    except PermissionError as pe:
        return {
            "status": "error",
            "message": str(pe)
        }
    except Exception as e:
        return {
            "status": "error",
            "message": f"An error occurred while searching in file: {str(e)}"
        }
